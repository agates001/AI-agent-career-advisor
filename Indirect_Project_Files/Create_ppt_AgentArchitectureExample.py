import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 1. Initialize Presentation with Widescreen Dimensions (16:9)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 2. Design Palettes (High-Contrast Corporate Dark/Cyber Aesthetic)
COLOR_BG = RGBColor(15, 23, 42)        # Slate 900
COLOR_PANEL = RGBColor(30, 41, 59)     # Slate 800
COLOR_TEXT = RGBColor(241, 245, 249)   # Slate 100
COLOR_MUTED = RGBColor(148, 163, 184)  # Slate 400
COLOR_ACCENT = RGBColor(56, 189, 248)  # Sky 400

# Helper function to apply dark theme background color
def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG

# Helper function to add structured headers
def add_slide_header(slide, title_text, subtitle_text):
    tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(1.2))
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p1 = tf.paragraphs[0]
    p1.text = title_text.upper()
    p1.font.name = "Arial"
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_ACCENT
    p1.space_after = Pt(4)
    
    p2 = tf.add_paragraph()
    p2.text = subtitle_text
    p2.font.name = "Arial"
    p2.font.size = Pt(14)
    p2.font.color.rgb = COLOR_MUTED

# ==============================================================================
# SLIDE 1: Title Slide
# ==============================================================================
slide_layout = prs.slide_layouts[6] # Blank Layout
slide1 = prs.slides.add_slide(slide_layout)
set_slide_background(slide1)

title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.333), Inches(3.0))
tf1 = title_box.text_frame
tf1.word_wrap = True

p1 = tf1.paragraphs[0]
p1.text = "THE ANATOMY OF A CREWAI AGENT OBJECT"
p1.font.name = "Arial"
p1.font.size = Pt(38)
p1.font.bold = True
p1.font.color.rgb = COLOR_ACCENT
p1.space_after = Pt(12)

p2 = tf1.add_paragraph()
p2.text = "Graduate Lecture Guide: Code Architecture, Component Subsystems, and Dynamic Execution Loops"
p2.font.name = "Arial"
p2.font.size = Pt(16)
p2.font.color.rgb = COLOR_TEXT

# ==============================================================================
# SLIDE 2: Core Components Component Map (Grid Layout)
# ==============================================================================
slide2 = prs.slides.add_slide(slide_layout)
set_slide_background(slide2)
add_slide_header(slide2, "Agent Object Structural Map", "The five explicit operational components configured inside the Python instantiation wrapper")

components = [
    {"name": "1. THE COGNITIVE BRAIN", "p": "llm=ollm", "desc": "Handles core token computation, text pattern analysis, semantic evaluations, and logical response parsing.", "libs": "openai, litellm"},
    {"name": "2. THE PERSONA ENGINE", "p": "role=, goal=, backstory=", "desc": "Compresses general knowledge boundaries into clear professional contexts using automated system prompt wrappers.", "libs": "crewai core"},
    {"name": "3. CONTEXTUAL MEMORY", "p": "memory=True", "desc": "Indexes localized execution context layers, runtime task handoffs, and learned entity facts.", "libs": "chromadb, pydantic"},
    {"name": "4. REAL-WORLD HANDS", "p": "tools=[...]", "desc": "Gives the agent API access hooks to scrape websites, query vector systems, read directories, or write code files.", "libs": "crewai-tools"},
    {"name": "5. REFLECTOR LOOP", "p": "verbose=True", "desc": "Governs the ReAct evaluation iteration step before outputting a validated solution block to the user.", "libs": "Python re / json"}
]

# Create a 5-column grid horizontally
left_margin = Inches(0.5)
top_margin = Inches(2.0)
col_width = Inches(2.2)
col_gap = Inches(0.2)
box_height = Inches(4.5)

for i, comp in enumerate(components):
    col_left = left_margin + i * (col_width + col_gap)
    # Add a panel backdrop shape
    shape = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, col_left, top_margin, col_width, box_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_PANEL
    shape.line.color.rgb = COLOR_MUTED
    shape.line.width = Pt(1)
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0.15)
    
    # Headline
    p = tf.paragraphs[0]
    p.text = comp["name"]
    p.font.name = "Arial"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT
    p.space_after = Pt(8)
    
    # Code Parameter
    p_code = tf.add_paragraph()
    p_code.text = f"Code: {comp['p']}"
    p_code.font.name = "Courier New"
    p_code.font.size = Pt(10)
    p_code.font.bold = True
    p_code.font.color.rgb = COLOR_TEXT
    p_code.space_after = Pt(12)
    
    # Description
    p_desc = tf.add_paragraph()
    p_desc.text = comp["desc"]
    p_desc.font.name = "Arial"
    p_desc.font.size = Pt(11)
    p_desc.font.color.rgb = COLOR_TEXT
    p_desc.space_after = Pt(20)
    
    # Libraries
    p_lib = tf.add_paragraph()
    p_lib.text = f"Under the Hood:\n{comp['libs']}"
    p_lib.font.name = "Arial"
    p_lib.font.size = Pt(10)
    p_lib.font.italic = True
    p_lib.font.color.rgb = COLOR_MUTED

# ==============================================================================
# SLIDE 3: The Dynamic Payload Assembly (Anatomy Matrix Layout)
# ==============================================================================
slide3 = prs.slides.add_slide(slide_layout)
set_slide_background(slide3)
add_slide_header(slide3, "The Dynamic Payload Assembly", "How the framework translates component variables into a singular prompt sent to the LLM API")

# Add a large combined text panel
panel = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.0), Inches(12.333), Inches(4.5))
panel.fill.solid()
panel.fill.fore_color.rgb = COLOR_PANEL
panel.line.color.rgb = COLOR_ACCENT
panel.line.width = Pt(1.5)

tf3 = panel.text_frame
tf3.word_wrap = True
tf3.margin_left = tf3.margin_top = tf3.margin_right = tf3.margin_bottom = Inches(0.25)

payload_layers = [
    ("[ PERSONA BLOCK ]", "Injects role, goal, and backstory constraints directly to skew the base alignment model."),
    ("[ MEMORY RETRIEVAL ]", "Stitches relevant semantic historical context fragments recovered from local vector database lookups."),
    ("[ TOOL REGISTRY ]", "Exposes the strict text structural arguments required to call interface hooks like file system readers."),
    ("[ CURRENT TASK BOUNDS ]", "Appends isolated operational goal scopes, specific criteria filters, and expected data schemas."),
    ("[ REACT INTERACTION TRIGGER ]", "Forces the cognitive runtime loop engine to open its internal evaluation step ('Thought:', 'Action:')")
]

for title, description in payload_layers:
    p_title = tf3.add_paragraph() if tf3.text else tf3.paragraphs[0]
    p_title.text = f"{title:<30} ──►   {description}"
    p_title.font.name = "Courier New"
    p_title.font.size = Pt(12)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_TEXT
    p_title.space_after = Pt(14)

# Save presentation output
prs.save("CrewAI_Agent_Anatomy_Lecture.pptx")
print("SUCCESS: 'CrewAI_Agent_Anatomy_Lecture.pptx' generated in your project folder.")
